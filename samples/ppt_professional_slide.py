"""Professional PowerPoint Slide — python-pptx로 전문가 수준 슬라이드 생성.

python-pptx로 정밀한 디자인 PPT 파일을 생성한 후,
AppleScript로 PowerPoint에서 자동으로 열어 실시간 시연합니다.

디자인 원칙:
- 그리드 기반 레이아웃 (16:9, 13.33 x 7.5 inch)
- 제한된 색상 팔레트 (4색)
- 일관된 폰트 체계 (타이틀→본문→캡션)
- 충분한 여백, 시각적 위계

Usage:
    python samples/ppt_professional_slide.py

macOS + Microsoft PowerPoint + python-pptx 필요.
pip install python-pptx
"""

import os
import subprocess
import time

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ═══════════════════════════════════════════════════════════
# Design System
# ═══════════════════════════════════════════════════════════
class C:
    """Color Palette."""
    NAVY = RGBColor(24, 42, 75)
    NAVY_LIGHT = RGBColor(35, 58, 95)
    BLUE = RGBColor(41, 98, 255)
    BLUE_LIGHT = RGBColor(230, 238, 255)
    ORANGE = RGBColor(255, 122, 0)
    ORANGE_LIGHT = RGBColor(255, 243, 230)
    GREEN = RGBColor(16, 185, 129)
    GREEN_LIGHT = RGBColor(230, 250, 243)
    PURPLE = RGBColor(124, 58, 237)
    PURPLE_LIGHT = RGBColor(245, 243, 255)
    WHITE = RGBColor(255, 255, 255)
    GRAY_50 = RGBColor(249, 250, 251)
    GRAY_100 = RGBColor(243, 244, 246)
    GRAY_300 = RGBColor(209, 213, 219)
    GRAY_500 = RGBColor(107, 114, 128)
    GRAY_700 = RGBColor(55, 65, 81)
    GRAY_900 = RGBColor(17, 24, 39)


def add_shape(slide, left, top, width, height, fill=None, line=None, line_width=Pt(0)):
    """도형 추가 헬퍼."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill=None):
    """원형 도형."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(1)):
    """모서리 둥근 사각형."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, size=12, color=C.GRAY_900, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font_name="Pretendard"):
    """텍스트 설정 헬퍼."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # 여백 설정
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    # 앵커
    try:
        tf.paragraphs[0].alignment = align
    except Exception:
        pass

    # 수직 정렬
    from pptx.oxml.ns import qn
    txBody = shape._element.txBody
    if txBody is not None:
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            anchor_map = {
                MSO_ANCHOR.TOP: 't',
                MSO_ANCHOR.MIDDLE: 'ctr',
                MSO_ANCHOR.BOTTOM: 'b',
            }
            bodyPr.set('anchor', anchor_map.get(anchor, 'ctr'))

    # 텍스트 (줄바꿈 지원)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name

    return tf


# ═══════════════════════════════════════════════════════════
# Slide 1: Title Page
# ═══════════════════════════════════════════════════════════
def create_slide_1(prs):
    """타이틀 페이지 — 좌측 네이비 사이드바 + 우측 메인."""
    print("  [Slide 1] 타이틀 페이지...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # ── 좌측 네이비 사이드바 ──
    add_shape(slide, Inches(0), Inches(0), Inches(4.2), Inches(7.5), fill=C.NAVY)

    # 장식 원 (반투명 느낌)
    add_oval(slide, Inches(-0.4), Inches(4.8), Inches(2.5), Inches(2.5), fill=C.NAVY_LIGHT)
    add_oval(slide, Inches(2.8), Inches(-0.3), Inches(1.0), Inches(1.0), fill=C.NAVY_LIGHT)

    # 로고
    logo = add_shape(slide, Inches(0.5), Inches(0.6), Inches(3), Inches(0.6))
    set_text(logo, "LogosAI", size=30, color=C.WHITE, bold=True)

    # 태그라인
    tag = add_shape(slide, Inches(0.5), Inches(1.2), Inches(3.2), Inches(0.4))
    set_text(tag, "Desktop Native AI Platform", size=12, color=RGBColor(150, 170, 210))

    # 구분선
    add_shape(slide, Inches(0.5), Inches(1.75), Inches(0.8), Inches(0.04), fill=C.BLUE)

    # 스펙 리스트
    specs = [
        "56+ AI Agents",
        "Desktop Native Control",
        "Self-Evolution System",
        "Enterprise Ready",
        "On-Premise LLM Support",
    ]
    for i, spec in enumerate(specs):
        s = add_shape(slide, Inches(0.5), Inches(2.1 + i * 0.38), Inches(3.2), Inches(0.35))
        set_text(s, spec, size=12, color=RGBColor(180, 195, 225))

    # ── 우측 메인 ──
    # 메인 타이틀
    title = add_shape(slide, Inches(4.7), Inches(1.2), Inches(8), Inches(1.0))
    set_text(title, "Agentic AI\nFramework", size=42, color=C.GRAY_900, bold=True,
             anchor=MSO_ANCHOR.BOTTOM)

    # 서브타이틀
    sub = add_shape(slide, Inches(4.7), Inches(2.4), Inches(7), Inches(0.7))
    set_text(sub, "에이전트가 스스로 학습하고, 진화하는\nDesktop Native AI 플랫폼",
             size=16, color=C.GRAY_500, anchor=MSO_ANCHOR.TOP)

    # 블루 구분선
    add_shape(slide, Inches(4.7), Inches(3.3), Inches(1.0), Inches(0.05), fill=C.BLUE)

    # 3개 지표 카드
    card_data = [
        ("56+", "AI Agents", C.BLUE, C.BLUE_LIGHT),
        ("L2+", "Autonomy", C.ORANGE, C.ORANGE_LIGHT),
        ("99.9%", "Uptime", C.GREEN, C.GREEN_LIGHT),
    ]

    for i, (num, label, accent, bg) in enumerate(card_data):
        x = Inches(4.7 + i * 2.6)
        y = Inches(3.7)
        w = Inches(2.3)
        h = Inches(1.5)

        # 카드 배경
        card = add_rounded_rect(slide, x, y, w, h, fill=bg, line=C.GRAY_100)

        # 상단 컬러 바
        add_shape(slide, x, y, w, Inches(0.05), fill=accent)

        # 숫자
        num_box = add_shape(slide, x, y + Inches(0.2), w, Inches(0.6))
        set_text(num_box, num, size=32, color=accent, bold=True, align=PP_ALIGN.CENTER)

        # 레이블
        lbl_box = add_shape(slide, x, y + Inches(0.85), w, Inches(0.4))
        set_text(lbl_box, label, size=12, color=C.GRAY_500, align=PP_ALIGN.CENTER)

    # 하단 바
    add_shape(slide, Inches(4.2), Inches(6.9), Inches(9.13), Inches(0.6), fill=C.GRAY_50)
    footer = add_shape(slide, Inches(4.5), Inches(7.0), Inches(4), Inches(0.3))
    set_text(footer, "2026 LogosAI — Confidential", size=9, color=C.GRAY_300)


# ═══════════════════════════════════════════════════════════
# Slide 2: Architecture
# ═══════════════════════════════════════════════════════════
def create_slide_2(prs):
    """아키텍처 다이어그램 — 3-Tier 카드."""
    print("  [Slide 2] 아키텍처 다이어그램...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 블루 바
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill=C.BLUE)

    # 타이틀
    t = add_shape(slide, Inches(0.7), Inches(0.3), Inches(5), Inches(0.5))
    set_text(t, "System Architecture", size=26, color=C.GRAY_900, bold=True)

    st = add_shape(slide, Inches(0.7), Inches(0.75), Inches(5), Inches(0.3))
    set_text(st, "3-Tier Desktop Native AI Architecture", size=13, color=C.GRAY_500)

    # 3개 Tier 카드
    tiers = [
        {
            "icon": "UI", "title": "Frontend", "sub": "logos_web · Next.js · Port 8010",
            "items": ["Chat UI + Artifact Preview", "SSE Streaming (14 Events)", "Conversation History"],
            "color": C.BLUE, "bg": C.BLUE_LIGHT,
        },
        {
            "icon": "API", "title": "Backend", "sub": "logos_api · FastAPI · Port 8090",
            "items": ["Orchestrator + Query Planner", "User Memory Engine", "Interaction Engine"],
            "color": C.ORANGE, "bg": C.ORANGE_LIGHT,
        },
        {
            "icon": "AI", "title": "Agent Runtime", "sub": "ACP Server · Port 8888",
            "items": ["56+ Specialized Agents", "ReAct + Tool Use + Memory", "Self-Evolution (L2+)"],
            "color": C.GREEN, "bg": C.GREEN_LIGHT,
        },
    ]

    tier_w = Inches(3.6)
    tier_h = Inches(3.8)
    tier_gap = Inches(0.4)
    start_x = Inches(0.7)
    tier_y = Inches(1.3)

    for i, tier in enumerate(tiers):
        x = start_x + i * (tier_w + tier_gap)

        # 카드 배경
        card = add_rounded_rect(slide, x, tier_y, tier_w, tier_h,
                                fill=C.WHITE, line=C.GRAY_100)

        # 상단 컬러 바
        add_shape(slide, x, tier_y, tier_w, Inches(0.06), fill=tier["color"])

        # 아이콘 원
        icon_size = Inches(0.65)
        icon_x = x + (tier_w - icon_size) // 2
        icon_y = tier_y + Inches(0.3)
        icon_shape = add_oval(slide, icon_x, icon_y, icon_size, icon_size, fill=tier["bg"])
        set_text(icon_shape, tier["icon"], size=16, color=tier["color"],
                 bold=True, align=PP_ALIGN.CENTER)

        # 타이틀
        title_box = add_shape(slide, x, tier_y + Inches(1.1), tier_w, Inches(0.4))
        set_text(title_box, tier["title"], size=20, color=C.GRAY_900,
                 bold=True, align=PP_ALIGN.CENTER)

        # 서브타이틀
        sub_box = add_shape(slide, x, tier_y + Inches(1.5), tier_w, Inches(0.3))
        set_text(sub_box, tier["sub"], size=10, color=C.GRAY_500, align=PP_ALIGN.CENTER)

        # 구분선
        add_shape(slide, x + Inches(0.5), tier_y + Inches(1.95),
                  tier_w - Inches(1.0), Inches(0.015), fill=C.GRAY_100)

        # 항목
        for j, item in enumerate(tier["items"]):
            item_y = tier_y + Inches(2.15 + j * 0.45)

            # 불릿 원
            add_oval(slide, x + Inches(0.4), item_y + Inches(0.08),
                     Inches(0.1), Inches(0.1), fill=tier["color"])

            # 텍스트
            item_box = add_shape(slide, x + Inches(0.65), item_y,
                                 tier_w - Inches(1.0), Inches(0.35))
            set_text(item_box, item, size=12, color=C.GRAY_700)

    # 화살표 (→)
    for i in range(2):
        ax = start_x + (i + 1) * tier_w + i * tier_gap + tier_gap // 2
        ay = tier_y + tier_h // 2 - Inches(0.15)
        arrow = add_shape(slide, ax - Inches(0.1), ay, tier_gap + Inches(0.2), Inches(0.3))
        set_text(arrow, "→", size=22, color=C.GRAY_300, align=PP_ALIGN.CENTER)

    # 하단 Desktop Native 배너
    banner_y = tier_y + tier_h + Inches(0.25)
    banner = add_rounded_rect(slide, start_x, banner_y,
                              tier_w * 3 + tier_gap * 2, Inches(0.55),
                              fill=C.PURPLE_LIGHT, line=RGBColor(230, 225, 255))
    set_text(banner,
             "Desktop Native:  PowerPoint · Excel · VS Code · KakaoTalk · Chrome — AI가 앱을 직접 제어",
             size=12, color=C.PURPLE, bold=True, align=PP_ALIGN.CENTER)

    # 페이지 번호
    pg = add_shape(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    set_text(pg, "2 / 3", size=9, color=C.GRAY_300, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
# Slide 3: Differentiation (Why LogosAI?)
# ═══════════════════════════════════════════════════════════
def create_slide_3(prs):
    """차별화 비교 — 2x2 피처 카드."""
    print("  [Slide 3] 차별화 비교...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 상단 바
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), fill=C.BLUE)

    # 타이틀
    t = add_shape(slide, Inches(0.7), Inches(0.3), Inches(5), Inches(0.5))
    set_text(t, "Why LogosAI?", size=26, color=C.GRAY_900, bold=True)

    st = add_shape(slide, Inches(0.7), Inches(0.75), Inches(6), Inches(0.3))
    set_text(st, "기존 AI 프레임워크와의 핵심 차별점", size=13, color=C.GRAY_500)

    # 2x2 카드
    features = [
        {
            "icon": "Native", "title": "Desktop Native Control",
            "desc": "PowerPoint, Excel, VS Code, KakaoTalk 등\n데스크톱 앱을 직접 제어합니다.\n웹 브라우저에 갇히지 않습니다.",
            "color": C.BLUE, "bg": C.BLUE_LIGHT,
        },
        {
            "icon": "Evolve", "title": "Self-Evolution (L2+)",
            "desc": "에이전트가 스스로 학습하고\n약한 부분을 자동 감지하여\nFORGE로 개선합니다.",
            "color": C.ORANGE, "bg": C.ORANGE_LIGHT,
        },
        {
            "icon": "Multi", "title": "Multi-Agent Orchestration",
            "desc": "56+ 전문 에이전트를\n직렬/병렬/하이브리드로\n자동 조합하여 실행합니다.",
            "color": C.GREEN, "bg": C.GREEN_LIGHT,
        },
        {
            "icon": "Secure", "title": "On-Premise Ready",
            "desc": "Ollama/vLLM으로\n폐쇄망 완전 구동 가능.\n금융/공공 규제 대응.",
            "color": C.PURPLE, "bg": C.PURPLE_LIGHT,
        },
    ]

    card_w = Inches(5.5)
    card_h = Inches(2.1)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    start_x = Inches(0.7)
    start_y = Inches(1.3)

    for i, feat in enumerate(features):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # 카드
        card = add_rounded_rect(slide, x, y, card_w, card_h,
                                fill=feat["bg"], line=C.GRAY_100)

        # 아이콘 원
        icon = add_oval(slide, x + Inches(0.3), y + (card_h - Inches(0.7)) // 2,
                        Inches(0.7), Inches(0.7), fill=feat["color"])
        set_text(icon, feat["icon"], size=11, color=C.WHITE, bold=True,
                 align=PP_ALIGN.CENTER)

        # 타이틀
        ttl = add_shape(slide, x + Inches(1.2), y + Inches(0.25),
                        card_w - Inches(1.5), Inches(0.35))
        set_text(ttl, feat["title"], size=17, color=C.GRAY_900, bold=True)

        # 설명
        desc = add_shape(slide, x + Inches(1.2), y + Inches(0.65),
                         card_w - Inches(1.5), Inches(1.2))
        set_text(desc, feat["desc"], size=12, color=C.GRAY_700,
                 anchor=MSO_ANCHOR.TOP)

    # CTA 버튼
    cta_y = Inches(5.9)
    cta = add_rounded_rect(slide, Inches(4.7), cta_y, Inches(3.8), Inches(0.55),
                           fill=C.BLUE)
    set_text(cta, "PoC 및 상세 데모 요청", size=15, color=C.WHITE, bold=True,
             align=PP_ALIGN.CENTER)

    # 페이지 번호
    pg = add_shape(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    set_text(pg, "3 / 3", size=9, color=C.GRAY_300, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════
def main():
    print("=" * 60)
    print("  Professional Slide Generator")
    print("  python-pptx + AppleScript 조합")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_slide_1(prs)
    create_slide_2(prs)
    create_slide_3(prs)

    # 저장
    output_path = os.path.expanduser("~/Desktop/LogosAI_Professional.pptx")
    prs.save(output_path)
    print(f"\n  -> 저장: {output_path}")

    # PowerPoint에서 열기
    print("  -> PowerPoint에서 열기...")
    subprocess.run(["open", "-a", "Microsoft PowerPoint", output_path])
    time.sleep(3)

    print("\n" + "=" * 60)
    print("  완료! 3장의 프로페셔널 슬라이드가 생성되었습니다.")
    print()
    print("  1. 타이틀 — 네이비 사이드바 + 지표 카드 3개")
    print("  2. 아키텍처 — 3-Tier 카드 + Desktop Native 배너")
    print("  3. 차별화 — 2x2 피처 카드 + CTA 버튼")
    print()
    print("  디자인:")
    print("  - 4색 팔레트 (Navy/Blue/Orange/Green)")
    print("  - 둥근 카드 + 아이콘 원 + 컬러 바")
    print("  - Pretendard 폰트 체계")
    print("  - 도형 70+ 개 프로그래밍 배치")
    print("=" * 60)


if __name__ == "__main__":
    main()
